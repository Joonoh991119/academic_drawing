#!/usr/bin/env python3
"""Robust PPTX extractor for reconstruction mode — deterministically pull every slide's title, body
text, speaker notes, tables, and autoshape inventory, AND save every embedded image INCLUDING pictures
embedded in PLACEHOLDERS.

A naive `shape_type == PICTURE` walk misses a figure embedded in a Content Placeholder (its shape_type
is PLACEHOLDER, not PICTURE) and silently drops it — so the rebuild can't re-embed what it never had,
and the loss is only caught post-hoc by reconstruction_diff.py. Per-run hand-written extractors catch
this inconsistently; this script makes it reliable by detecting images via `shape.image` on EVERY
shape (recursing groups), not by shape type.

    python3 pptx_extract.py DECK.pptx --out OUTDIR

Writes OUTDIR/00_extracted.md and OUTDIR/figs/slideNN_imgMM.<ext>. READ-ONLY on the source.
"""
import argparse
import sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
except Exception as e:  # pragma: no cover
    print(f"[pptx-extract] python-pptx required: {e}", file=sys.stderr)
    sys.exit(3)

_A = "{http://schemas.openxmlformats.org/drawingml/2006/main}"


def _para_lines(tf):
    lines = []
    for p in tf.paragraphs:
        t = ("".join(r.text for r in p.runs) or p.text or "").strip()
        if t:
            lines.append("  " * (p.level or 0) + t)
    return lines


def extract(path, out):
    prs = Presentation(path)
    figs = out / "figs"
    figs.mkdir(parents=True, exist_ok=True)
    md = [f"# Extracted: {Path(path).name}", "", "", ""]  # md[2] reserved for the summary line
    n_imgs = n_tables = 0

    for si, slide in enumerate(prs.slides, 1):
        try:
            title_shape = slide.shapes.title
        except Exception:
            title_shape = None
        title = None
        if title_shape is not None:
            try:
                title = (title_shape.text or "").strip() or None
            except Exception:
                title = None

        body, tables, autoshapes = [], [], []
        state = {"img_idx": 0}

        def walk(shapes):
            for sh in shapes:
                try:
                    st = sh.shape_type
                except Exception:
                    st = None
                if st == MSO_SHAPE_TYPE.GROUP:
                    walk(sh.shapes)
                    continue
                # IMAGE — any shape carrying a blob (PICTURE or placeholder-embedded picture)
                try:
                    blob, ext = sh.image.blob, sh.image.ext
                except Exception:
                    blob, ext = None, None
                if blob is not None:
                    state["img_idx"] += 1
                    nonlocal_n = state["img_idx"]
                    (figs / f"slide{si:02d}_img{nonlocal_n:02d}.{ext or 'png'}").write_bytes(blob)
                    continue
                # TABLE
                if getattr(sh, "has_table", False):
                    rows = [" | ".join(c.text.strip() for c in r.cells) for r in sh.table.rows]
                    tables.append(rows)
                    continue
                # AUTOSHAPE geometry preset (potential visual-argument carrier: arrow / no-entry / bracket)
                try:
                    spPr = sh._element.spPr
                    geom = spPr.find(f"{_A}prstGeom") if spPr is not None else None
                    if geom is not None:
                        preset = geom.get("prst")
                        if preset and preset != "rect":
                            autoshapes.append(preset)
                except Exception:
                    pass
                # BODY TEXT (skip the title shape — already captured)
                if sh is not title_shape and getattr(sh, "has_text_frame", False) and sh.has_text_frame:
                    body.extend(_para_lines(sh.text_frame))

        walk(slide.shapes)
        n_imgs += state["img_idx"]
        n_tables += len(tables)

        md.append(f"## Slide {si}")
        md.append(f"**Title:** {title or '(none)'}")
        md.append(f"**Images:** {state['img_idx']}")
        if body:
            md.append("**Body:**")
            md.extend(f"- {b}" for b in body)
        if tables:
            md.append("**Tables:**")
            for t in tables:
                md.extend(f"  {row}" for row in t)
        if autoshapes:
            md.append(f"**Autoshapes (prstGeom, non-rect — may carry the visual argument):** "
                      f"{', '.join(autoshapes)}")
        try:
            if slide.has_notes_slide:
                note = (slide.notes_slide.notes_text_frame.text or "").strip()
                if note:
                    md.append(f"**Notes:** {note}")
        except Exception:
            pass
        md.append("")

    md[2] = f"_slides={len(prs.slides)} · images_saved={n_imgs} · tables={n_tables}_"
    (out / "00_extracted.md").write_text("\n".join(md))
    print(f"[pptx-extract] slides={len(prs.slides)} images={n_imgs} tables={n_tables} "
          f"-> {out}/00_extracted.md (+ figs/)")


def main():
    ap = argparse.ArgumentParser(description="Robust reconstruction-mode PPTX extractor")
    ap.add_argument("deck")
    ap.add_argument("--out", required=True)
    a = ap.parse_args()
    try:
        extract(a.deck, Path(a.out))
    except Exception as e:
        print(f"[pptx-extract] failed: {e}", file=sys.stderr)
        sys.exit(3)


if __name__ == "__main__":
    main()
