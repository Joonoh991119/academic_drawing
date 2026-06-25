#!/usr/bin/env python3
"""Reconstruction FIDELITY DIFF (deck-rebuild mode) — make silent content drops VISIBLE.

When the harness reconstructs an existing deck, the one-shot builder tends to "improve" past
fidelity: re-author the slide count, collapse a dense table to prose, or drop a real result figure.
slide-builder's rules forbid this, but rules alone keep getting violated — so this is the
deterministic gate that COUNTS what changed and, crucially, MATCHES each source figure by content
hash so a drop-and-swap (counts unchanged but a figure replaced) can't slip through.

    python3 reconstruction_diff.py ORIGINAL.pptx RECON.pptx [--json out.json] [--strict]

Reports, for both decks: slide count, image-bearing-shape count, table count (recursing groups). Then
hashes every embedded image blob and reports which SOURCE figures have no byte-identical match in the
reconstruction. Everything is WARN, not FAIL, on purpose: this script cannot tell a DATA figure (must
keep) from a DECORATIVE icon (OK to drop), nor a table rebuilt as a shape-grid (content kept) from one
collapsed to prose (content lost), nor a figure legitimately TRANSFORMED (WMF->PNG, cropped, recolored
— hash changes though content is preserved) from one truly dropped. Its job is to ensure no change is
SILENT; each delta and each unmatched figure must be DISPOSITIONED in the orig->recon manifest.
`--strict` turns any delta/unmatched figure into a nonzero exit so a pipeline can block until the
manifest + confirmation exist.

Exit: 0 = no deltas (or non-strict run), 2 = deltas present under --strict, 3 = could not read a deck.
"""
import argparse
import hashlib
import json
import sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
except Exception as e:  # pragma: no cover
    print(f"[recon-diff] python-pptx required: {e}", file=sys.stderr)
    sys.exit(3)


def _inventory(path):
    """Count slides, image-bearing shapes, tables, and collect image-blob hashes.

    Counts ANY shape carrying an embedded image (a PICTURE shape OR a picture placeholder), not just
    shape_type==PICTURE — a result figure stored in a placeholder would otherwise be undercounted and a
    real drop could be missed. Each image's blob is sha1-hashed so figures can be matched by identity
    across the two decks (the rebuild re-embeds the same bytes via addImage, so a kept figure's hash
    matches; a dropped one's does not).
    """
    prs = Presentation(path)
    pics = tables = 0
    hashes = set()

    def walk(shapes):
        nonlocal pics, tables
        for sh in shapes:
            try:
                st = sh.shape_type
            except Exception:
                st = None
            if st == MSO_SHAPE_TYPE.GROUP:
                walk(sh.shapes)
                continue
            blob = None
            try:
                blob = sh.image.blob
            except Exception:
                blob = None
            if blob is not None:
                pics += 1
                hashes.add(hashlib.sha1(blob).hexdigest())
            if getattr(sh, "has_table", False):
                tables += 1

    for sl in prs.slides:
        walk(sl.shapes)
    return {"slides": len(prs.slides), "pictures": pics, "tables": tables, "_hashes": hashes}


def diff(original, recon):
    o = _inventory(original)
    r = _inventory(recon)
    findings = []
    dp = o["pictures"] - r["pictures"]
    dt = o["tables"] - r["tables"]
    ds = r["slides"] - o["slides"]
    unmatched = len(o["_hashes"] - r["_hashes"])  # unique source figures with no byte match in recon
    if dp > 0:
        findings.append({
            "severity": "WARN", "kind": "figures-fewer",
            "detail": f"{dp} fewer figure(s) ({o['pictures']}->{r['pictures']}). Confirm EVERY dropped "
                      f"image is decorative — a real data/result figure must be re-embedded, not lost.",
        })
    if unmatched > 0:
        findings.append({
            "severity": "WARN", "kind": "figures-unmatched",
            "detail": f"{unmatched} source figure(s) have NO byte-identical match in the reconstruction "
                      f"(dropped, OR legitimately transformed: cropped / WMF->PNG / recolored / "
                      f"re-typeset). Account for each — a count match alone can hide a drop-and-swap.",
        })
    if dt > 0:
        findings.append({
            "severity": "WARN", "kind": "tables-fewer",
            "detail": f"{dt} fewer table(s) ({o['tables']}->{r['tables']}). Confirm each was rebuilt as "
                      f"a native/shape grid, NOT collapsed to prose or summary scores.",
        })
    if ds != 0:
        findings.append({
            "severity": "WARN", "kind": "slide-count-changed",
            "detail": f"slide count {o['slides']}->{r['slides']} (delta {ds:+d}). Requires an orig->recon "
                      f"manifest (kept/merged/dropped/added, each justified).",
        })
    pub = lambda d: {"slides": d["slides"], "pictures": d["pictures"], "tables": d["tables"]}
    return {
        "original": pub(o), "recon": pub(r),
        "figures_unmatched": unmatched, "source_figures_unique": len(o["_hashes"]),
        "verdict": "WARN" if findings else "PASS", "warn": len(findings), "findings": findings,
    }


def main():
    ap = argparse.ArgumentParser(description="Deterministic reconstruction fidelity diff")
    ap.add_argument("original")
    ap.add_argument("recon")
    ap.add_argument("--json", default=None)
    ap.add_argument("--strict", action="store_true", help="nonzero exit when any delta exists")
    a = ap.parse_args()
    try:
        rep = diff(a.original, a.recon)
    except Exception as e:
        print(f"[recon-diff] could not read decks: {e}", file=sys.stderr)
        sys.exit(3)
    if a.json:
        Path(a.json).write_text(json.dumps(rep, indent=2))
    o, r = rep["original"], rep["recon"]
    print(f"[recon-diff] {rep['verdict']}  slides {o['slides']}->{r['slides']}  "
          f"pics {o['pictures']}->{r['pictures']}  tables {o['tables']}->{r['tables']}  "
          f"unmatched-figs {rep['figures_unmatched']}/{rep['source_figures_unique']}  WARN={rep['warn']}")
    for f in rep["findings"]:
        print(f"  [{f['severity']}] {f['kind']}: {f['detail']}")
    sys.exit(2 if (a.strict and rep["findings"]) else 0)


if __name__ == "__main__":
    main()
