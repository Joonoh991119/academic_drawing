#!/usr/bin/env python3
"""Deterministic PALETTE / LABEL-COLOR / hue-budget / FONT-minima gate for .pptx decks.

This is the deck's twin of overlap_check.py. Where overlap_check measures rendered SVG geometry
to catch text/shape collisions, this script walks a PowerPoint deck's XML and checks that it obeys
the project STYLE CONTRACT defined in palette.json (the single source of truth):

  FAIL  a STRUCTURAL fill hex is NOT one of the structural token hexes
  FAIL  a run whose text matches a label_map key does NOT carry that token's hex
  FAIL  a slide uses more distinct structural hues than palette.max_colors_per_slide
  FAIL  a font run is smaller than typography.font_minima_pt.slide_body / slide_title
  WARN  a fill/font color the file does not resolve to a literal RGB (theme color, scheme color,
        "auto", inherited-from-placeholder) — these are handed to the vision pass, mirroring
        overlap_check's exit-3 graceful-degrade philosophy: we never silently pass an unknown.

python-pptx exposes only what is written explicitly on a run/shape; colors that come from the
theme or a slide-layout placeholder are reported as theme/None and become WARNs, not FAILs — a
deterministic gate must not assert a violation it cannot actually see.

Usage:   python3 pptx_style_lint.py INPUT.pptx [--palette palette.json] [--json out.json] [--strict]
Exit:    0 = clean (no FAIL),  2 = FAIL present,  3 = could not run (missing dep/file/data),
         4 = WARN present under --strict
"""
import argparse, json, sys
from pathlib import Path

# --- dependency + file graceful-degrade (exit 3, never silently pass) ----------------------------
try:
    from pptx import Presentation
    from pptx.util import Emu
    from pptx.enum.dml import MSO_THEME_COLOR, MSO_FILL
    from pptx.dml.color import RGBColor
except Exception as e:  # python-pptx missing or broken
    print(f"[pptx-lint] python-pptx not importable ({e}); cannot run. "
          "Fallback: render the deck to images and use the vision pass.", file=sys.stderr)
    sys.exit(3)

# Default palette location relative to this script (…/overlap-qc/scripts/ -> …/ga-style-contract/assets/)
DEFAULT_PALETTE = (Path(__file__).resolve().parent.parent.parent
                   / "ga-style-contract" / "assets" / "palette.json")


def norm_hex(h):
    """Normalise a hex string to uppercase 6-digit '#RRGGBB' for set membership comparison."""
    if h is None:
        return None
    h = str(h).strip().lstrip("#").upper()
    if len(h) == 3:  # expand shorthand #abc -> #aabbcc
        h = "".join(c * 2 for c in h)
    if len(h) != 6 or any(c not in "0123456789ABCDEF" for c in h):
        return None
    return "#" + h


def _project_label_map():
    """Project condition->token map the Director writes to _workspace/00_input/label_map.json."""
    import os
    d = Path(os.getcwd())
    for _ in range(4):
        f = d / "_workspace" / "00_input" / "label_map.json"
        if f.exists():
            try:
                return {k: v for k, v in json.loads(f.read_text()).items() if not k.startswith("_")}
            except Exception:
                return {}
        d = d.parent
    return {}


def load_palette(path):
    """Load palette.json and derive the contract sets/values. Returns a dict or raises."""
    pal = json.loads(Path(path).read_text())
    structural = pal.get("structural", {})
    # structural token hexes (skip private "_desc"-style keys that have no hex)
    tokens = {name: norm_hex(v["hex"])
              for name, v in structural.items()
              if isinstance(v, dict) and "hex" in v}
    hex_set = {h for h in tokens.values() if h}
    # name(token) lookup for nicer findings, e.g. "#3D6E9B" -> "cond_a"
    hex_to_name = {h: name for name, h in tokens.items() if h}

    # paper/bg/text do NOT count toward the per-slide hue budget (per palette _desc)
    non_counting = {tokens.get(t) for t in ("paper", "bg", "text") if tokens.get(t)}

    # label_map: palette real keys + the project map (_workspace/00_input/label_map.json),
    # each resolved label -> required hex via its token name.
    merged = _project_label_map() or {k: v for k, v in pal.get("label_map", {}).items() if not k.startswith("_")}
    label_map = {}
    for label, token in merged.items():
        if not isinstance(token, str):
            continue
        req = tokens.get(token)
        if req:
            label_map[label.strip().lower()] = {"token": token, "hex": req}

    typ = pal.get("typography", {})
    minima = typ.get("font_minima_pt", {})
    return {
        "tokens": tokens,
        "hex_set": hex_set,
        "hex_to_name": hex_to_name,
        "non_counting": non_counting,
        "label_map": label_map,
        "max_per_slide": int(pal.get("max_colors_per_slide", 5)),
        "min_body": float(minima.get("slide_body", 0) or 0),
        "min_title": float(minima.get("slide_title", 0) or 0),
    }


def resolve_color(color_format):
    """Return ('rgb', '#RRGGBB') for a literal RGB fill/font color, ('theme', name) for a
    theme/scheme color, or ('none', None) when nothing is explicitly set / cannot be resolved.
    python-pptx raises if you read .rgb on a non-RGB color, so we branch on .type first."""
    try:
        ctype = color_format.type
    except Exception:
        return ("none", None)
    if ctype is None:
        return ("none", None)
    # MSO_THEME_COLOR member => theme color (resolved from the theme at render time, not here)
    if ctype == 1 or str(ctype).startswith("RGB") or getattr(ctype, "name", "") == "RGB":
        try:
            return ("rgb", norm_hex(str(color_format.rgb)))
        except Exception:
            return ("none", None)
    # anything else (THEME / scheme) is not a literal RGB we can verify deterministically
    try:
        return ("theme", getattr(color_format.theme_color, "name", str(color_format.theme_color)))
    except Exception:
        return ("theme", str(ctype))


def shape_fill_color(shape):
    """Resolve a shape's solid-fill color. Returns (kind, value) like resolve_color, or
    ('skip', None) when the shape has no fill / a background fill we should ignore."""
    try:
        fill = shape.fill
    except Exception:
        return ("skip", None)
    try:
        ftype = fill.type
    except Exception:
        return ("skip", None)
    if ftype is None:
        # inherited fill: not written on the shape, defer to vision
        return ("none", None)
    if ftype != MSO_FILL.SOLID:
        # gradient / picture / pattern / background -> not a single structural hue we gate here
        return ("skip", None)
    return resolve_color(fill.fore_color)


def is_title_placeholder(shape):
    """True if the shape is a title/centered-title placeholder (uses the title minimum)."""
    try:
        if not shape.is_placeholder:
            return False
        ph = shape.placeholder_format.type
        return ph is not None and "TITLE" in str(ph)
    except Exception:
        return False


def iter_runs(shape):
    """Yield (run, paragraph) for every text run in a shape that has a text frame."""
    if not getattr(shape, "has_text_frame", False):
        return
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            yield run, para


def lint(prs, P):
    """Walk the deck and return the report dict {verdict, fail, warn, findings, ...}."""
    findings = []

    # palette self-check: each structural token carries ONE meaning (no label_map overloading) —
    # 'accent' overloaded = FAIL (it must mark the single finding only, never a condition color).
    tok_labels = {}
    for lab, meta in P.get("label_map", {}).items():
        tok_labels.setdefault(meta["token"], []).append(lab)
    for tok, labs in tok_labels.items():
        if len(labs) > 1:
            findings.append({
                "severity": "FAIL" if tok == "accent" else "WARN", "kind": "token-overload",
                "slide": 0, "where": "label_map", "hex": None,
                "detail": (f"token '{tok}' is assigned to {len(labs)} labels ({', '.join(labs)}); "
                           + ("'accent' must mark the single finding only — never a condition color"
                              if tok == "accent" else "each token should carry one meaning project-wide"))})

    for s_idx, slide in enumerate(prs.slides, start=1):
        slide_hues = set()  # distinct structural RGB hexes seen on this slide (for the budget)

        for shape in slide.shapes:
            shp_name = shape.name or shape.shape_type
            is_title = is_title_placeholder(shape)

            # --- structural fill check ---------------------------------------------------------
            kind, val = shape_fill_color(shape)
            if kind == "rgb" and val:
                if val not in P["non_counting"]:
                    slide_hues.add(val)
                if val not in P["hex_set"]:
                    findings.append({
                        "severity": "FAIL", "kind": "off-palette-fill", "slide": s_idx,
                        "where": shp_name, "hex": val,
                        "detail": f"shape fill {val} is not a structural token "
                                  f"(allowed: {', '.join(sorted(P['hex_set']))})"})
            elif kind == "theme":
                findings.append({
                    "severity": "WARN", "kind": "theme-fill", "slide": s_idx,
                    "where": shp_name, "hex": None,
                    "detail": f"fill is a theme/scheme color ({val}); cannot verify against palette "
                              "deterministically — confirm via vision"})
            # kind in ("none","skip") -> nothing explicit to gate

            # --- per-run font checks -----------------------------------------------------------
            for run, _para in iter_runs(shape):
                text = (run.text or "").strip()
                font = run.font

                # font-size minimum (title vs body)
                size = font.size
                if size is not None:
                    pt = size.pt
                    floor = P["min_title"] if is_title else P["min_body"]
                    which = "slide_title" if is_title else "slide_body"
                    if floor and pt < floor:
                        findings.append({
                            "severity": "FAIL", "kind": "font-too-small", "slide": s_idx,
                            "where": shp_name, "text": text[:40], "pt": round(pt, 1),
                            "detail": f"{pt:.1f}pt < {which} minimum {floor:g}pt"})

                # run font color: feeds label check + budget + off-palette detection
                fkind, fval = resolve_color(font.color)
                if fkind == "rgb" and fval and fval not in P["non_counting"]:
                    slide_hues.add(fval)

                # label-map color binding: if a run's text IS a contracted label, its color
                # must be that label's token hex.
                key = text.lower()
                if key in P["label_map"]:
                    want = P["label_map"][key]
                    if fkind == "rgb" and fval:
                        if fval != want["hex"]:
                            findings.append({
                                "severity": "FAIL", "kind": "label-color-mismatch", "slide": s_idx,
                                "where": shp_name, "text": text[:40], "hex": fval,
                                "detail": f"label '{text}' must use token '{want['token']}' "
                                          f"({want['hex']}) but is {fval}"})
                    else:
                        # color not explicitly set / theme -> can't confirm the binding
                        findings.append({
                            "severity": "WARN", "kind": "label-color-unresolved", "slide": s_idx,
                            "where": shp_name, "text": text[:40], "hex": None,
                            "detail": f"label '{text}' should carry token '{want['token']}' "
                                      f"({want['hex']}) but its run color is {fkind} (unresolved) "
                                      "— confirm via vision"})

        # --- per-slide hue budget --------------------------------------------------------------
        if len(slide_hues) > P["max_per_slide"]:
            findings.append({
                "severity": "FAIL", "kind": "too-many-hues", "slide": s_idx,
                "where": f"slide {s_idx}", "n": len(slide_hues),
                "detail": f"{len(slide_hues)} distinct structural hues "
                          f"({', '.join(sorted(slide_hues))}) exceed max_colors_per_slide="
                          f"{P['max_per_slide']} (paper/bg/text excluded)"})

    fails = [f for f in findings if f["severity"] == "FAIL"]
    warns = [f for f in findings if f["severity"] == "WARN"]
    return {
        "n_slides": len(prs.slides),
        "fail": len(fails), "warn": len(warns),
        "verdict": "FAIL" if fails else ("WARN" if warns else "PASS"),
        "findings": findings,
    }


def main():
    ap = argparse.ArgumentParser(description="Deterministic .pptx palette/label/hue/font gate.")
    ap.add_argument("input", help="path to the .pptx deck")
    ap.add_argument("--palette", default=str(DEFAULT_PALETTE), help="path to palette.json")
    ap.add_argument("--json", default=None, help="report output path (default: INPUT.style.json)")
    ap.add_argument("--strict", action="store_true", help="WARNs also fail the gate (exit 4)")
    args = ap.parse_args()

    in_path = Path(args.input)
    if not in_path.exists():
        print(f"[pptx-lint] input not found: {in_path}", file=sys.stderr); sys.exit(3)
    pal_path = Path(args.palette)
    if not pal_path.exists():
        print(f"[pptx-lint] palette.json not found: {pal_path}; cannot define the contract.",
              file=sys.stderr); sys.exit(3)
    try:
        P = load_palette(pal_path)
    except Exception as e:
        print(f"[pptx-lint] could not parse palette.json ({e}).", file=sys.stderr); sys.exit(3)
    if not P["hex_set"]:
        print("[pptx-lint] palette.json has no structural token hexes; nothing to gate against.",
              file=sys.stderr); sys.exit(3)
    try:
        prs = Presentation(str(in_path))
    except Exception as e:
        print(f"[pptx-lint] could not open deck ({e}). Fallback: render to images + vision pass.",
              file=sys.stderr); sys.exit(3)

    rep = lint(prs, P)
    rep["input"] = str(in_path.resolve())
    rep["palette"] = str(pal_path.resolve())
    out = Path(args.json) if args.json else in_path.with_suffix(".style.json")
    out.write_text(json.dumps(rep, indent=2))

    print(f"[pptx-lint] {rep['verdict']}  slides={rep['n_slides']} "
          f"FAIL={rep['fail']} WARN={rep['warn']}  -> {out}")
    for f in rep["findings"]:
        loc = f.get("text") or f.get("hex") or f.get("where", "")
        print(f"  [{f['severity']}] s{f['slide']} {f['kind']}: {loc}  — {f['detail']}")

    if rep["fail"] > 0:
        sys.exit(2)
    if args.strict and rep["warn"] > 0:
        sys.exit(4)
    sys.exit(0)


if __name__ == "__main__":
    main()
